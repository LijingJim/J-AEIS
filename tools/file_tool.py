"""
通用文件解析工具。
目标是先识别文件类型，再尽可能提取可供 LLM 阅读的文本/结构摘要。
"""

import io
import json
import threading
import zipfile
from pathlib import Path

import pandas as pd
import streamlit as st
from PIL import Image

from tools.excel_tool import _load_df, _FILE_READ_LOCK, _get_cached_file_data, _set_cached_file_data


TEXT_EXTS = {
    ".txt", ".md", ".log", ".py", ".js", ".ts", ".json", ".xml", ".html",
    ".htm", ".csv", ".sql", ".yaml", ".yml", ".ini", ".cfg",
}
IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".gif", ".bmp", ".webp", ".tif", ".tiff"}
EXCEL_EXTS = {".xls", ".xlsx", ".csv"}


def _uploaded_file():
    """读取上传文件。优先从模块级缓存读（跨线程可见）。"""
    cached_bytes, cached_name = _get_cached_file_data()
    if cached_bytes is not None:
        cached_file = type("_CachedFile", (), {"name": cached_name})()
        return cached_file, cached_bytes

    file = st.session_state.get("uploaded_file")
    if file is None:
        raise ValueError("请先上传文件。")

    if hasattr(file, "closed") and file.closed:
        raise ValueError("文件对象已关闭，请重新上传文件。")

    with _FILE_READ_LOCK:
        try:
            file.seek(0)
            data = file.read()
            file.seek(0)
        except ValueError as e:
            raise ValueError(f"读取文件失败（文件可能已关闭或损坏）: {e}")
    if not data:
        raise ValueError("文件读取为空（0 字节），请确认文件内容完整。")
    _set_cached_file_data(data, file.name)
    return file, data


def _decode_text(data: bytes) -> str:
    for encoding in ("utf-8-sig", "utf-8", "gb18030", "gbk", "latin-1"):
        try:
            return data.decode(encoding)
        except UnicodeDecodeError:
            continue
    return data.decode("utf-8", errors="replace")


def _read_pdf(data: bytes, max_pages: int) -> str:
    try:
        from pypdf import PdfReader
    except ImportError:
        return "PDF 解析依赖 pypdf 未安装，请运行：pip install pypdf"

    reader = PdfReader(io.BytesIO(data))
    pages = []
    for i, page in enumerate(reader.pages[:max_pages], start=1):
        text = page.extract_text() or ""
        pages.append(f"## 第 {i} 页\n{text.strip()}")
    return f"PDF 共 {len(reader.pages)} 页，已提取前 {min(max_pages, len(reader.pages))} 页：\n\n" + "\n\n".join(pages)


def _read_docx(data: bytes, max_chars: int) -> str:
    try:
        from docx import Document
    except ImportError:
        return "Word 解析依赖 python-docx 未安装，请运行：pip install python-docx"

    doc = Document(io.BytesIO(data))
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    text = "\n".join(paragraphs)
    return text[:max_chars] + ("..." if len(text) > max_chars else "")


def _read_pptx(data: bytes, max_chars: int) -> str:
    try:
        from pptx import Presentation
    except ImportError:
        return "PPT 解析依赖 python-pptx 未安装，请运行：pip install python-pptx"

    prs = Presentation(io.BytesIO(data))
    chunks = []
    for idx, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())
        if texts:
            chunks.append(f"## 第 {idx} 页\n" + "\n".join(texts))
    text = "\n\n".join(chunks)
    return text[:max_chars] + ("..." if len(text) > max_chars else "")


def inspect_uploaded_file(max_chars: int = 6000, max_pages: int = 5) -> str:
    """
    通用文件检查/读取。
    支持：Excel/CSV、文本/代码、JSON、PDF、Word、PPT、HTML表格、图片、ZIP。
    """
    try:
        file, data = _uploaded_file()
        name = file.name
        ext = Path(name).suffix.lower()
        size_mb = len(data) / 1024 / 1024

        header = [f"文件名：{name}", f"大小：{size_mb:.2f} MB", f"扩展名：{ext or '无'}", ""]

        if ext in EXCEL_EXTS:
            df = _load_df(file, nrows=50)
            return "\n".join(header) + f"表格预览：{len(df)} 行 × {len(df.columns)} 列\n列名：{list(df.columns)}\n\n{df.to_markdown(index=False)}"

        if ext == ".json":
            text = _decode_text(data)
            obj = json.loads(text)
            pretty = json.dumps(obj, ensure_ascii=False, indent=2)
            return "\n".join(header) + pretty[:max_chars] + ("..." if len(pretty) > max_chars else "")

        if ext in {".html", ".htm"}:
            try:
                tables = pd.read_html(io.BytesIO(data))
                if tables:
                    df = tables[0].head(50)
                    return "\n".join(header) + f"检测到 {len(tables)} 个表格，预览第 1 个：\n\n{df.to_markdown(index=False)}"
            except Exception:
                pass
            text = _decode_text(data)
            return "\n".join(header) + text[:max_chars] + ("..." if len(text) > max_chars else "")

        if ext == ".pdf":
            return "\n".join(header) + _read_pdf(data, max_pages)

        if ext == ".docx":
            return "\n".join(header) + _read_docx(data, max_chars)

        if ext == ".pptx":
            return "\n".join(header) + _read_pptx(data, max_chars)

        if ext in IMAGE_EXTS:
            image = Image.open(io.BytesIO(data))
            st.session_state["uploaded_image_preview"] = image.copy()
            return "\n".join(header) + (
                f"图片信息：格式={image.format}，尺寸={image.width}×{image.height}，模式={image.mode}\n\n"
                "图片已在界面中预览。注意：当前 deepseek-chat 是文本模型，不能直接理解图片内容；"
                "如果需要识别图片文字或图像内容，需要接入 OCR 或视觉模型。"
            )

        if ext == ".zip":
            with zipfile.ZipFile(io.BytesIO(data)) as zf:
                names = zf.namelist()
            preview = "\n".join(f"- {n}" for n in names[:100])
            return "\n".join(header) + f"ZIP 内共 {len(names)} 个文件，前 100 个：\n{preview}"

        if ext in TEXT_EXTS or data[:1024].count(b"\x00") == 0:
            text = _decode_text(data)
            return "\n".join(header) + text[:max_chars] + ("..." if len(text) > max_chars else "")

        return "\n".join(header) + "暂无法提取该二进制文件内容。可先转换为文本、表格、PDF、图片或压缩包后再上传。"
    except Exception as e:
        return f"文件解析失败：{e}"


# ── 长文档工具链（PDF 深度支持）─────────────────

def list_pdf_info() -> str:
    """返回 PDF 文件元信息：总页数、标题、作者、是否加密等。用于长文档分析第一步——了解文档规模。"""
    try:
        from pypdf import PdfReader
    except ImportError:
        return "PDF 解析依赖 pypdf 未安装，请运行：pip install pypdf"

    try:
        _, data = _uploaded_file()
        reader = PdfReader(io.BytesIO(data))
        meta = reader.metadata or {}
        info_lines = [
            f"PDF 文件信息：",
            f"- 总页数：{len(reader.pages)}",
            f"- 标题：{getattr(meta, 'title', None) or meta.get('/Title', '未知')}",
            f"- 作者：{getattr(meta, 'author', None) or meta.get('/Author', '未知')}",
            f"- 主题：{getattr(meta, 'subject', None) or meta.get('/Subject', '未知')}",
            f"- 是否加密：{'是' if reader.is_encrypted else '否'}",
        ]
        # 尝试提取目录
        toc = _extract_toc_from_reader(reader)
        if toc:
            info_lines.append(f"\n📑 文档目录结构（{len(toc)} 条）：")
            for entry in toc[:30]:
                info_lines.append(f"  [{entry['page']}] {entry['title']}")
            if len(toc) > 30:
                info_lines.append(f"  ... 共 {len(toc)} 条，以上为前 30 条")
        else:
            info_lines.append("\n⚠️ 该 PDF 无内嵌目录，建议用 search_pdf 按关键词定位内容。")
        return "\n".join(info_lines)
    except Exception as e:
        return f"获取 PDF 信息失败：{e}"


def _extract_toc_from_reader(reader) -> list:
    """从 PdfReader 提取目录结构（递归展平）。"""
    try:
        from pypdf import PdfReader
    except ImportError:
        return []

    try:
        # pypdf 4.x+ 使用 reader.outline
        if hasattr(reader, 'outline') and reader.outline:
            result = []

            def _flatten(items, level=0):
                for item in items:
                    if isinstance(item, list):
                        _flatten(item, level + 1)
                    elif hasattr(item, 'title'):
                        page_num = None
                        if hasattr(item, 'page') and hasattr(item.page, 'page_number'):
                            page_num = item.page.page_number
                        elif hasattr(item, 'page_number'):
                            page_num = item.page_number
                        else:
                            try:
                                page_num = reader.get_page_number(item) + 1
                            except Exception:
                                page_num = "?"
                        prefix = "  " * level + ("└─ " if level > 0 else "")
                        result.append({"title": f"{prefix}{item.title}", "page": page_num})

            _flatten(reader.outline)
            return result
    except Exception:
        pass
    return []


def extract_pdf_toc() -> str:
    """提取 PDF 的目录（书签/大纲）结构。如果 PDF 有内嵌书签，可快速了解文档章节分布。"""
    try:
        from pypdf import PdfReader
    except ImportError:
        return "PDF 解析依赖 pypdf 未安装，请运行：pip install pypdf"

    try:
        _, data = _uploaded_file()
        reader = PdfReader(io.BytesIO(data))
        toc = _extract_toc_from_reader(reader)
        if not toc:
            return "该 PDF 无内嵌目录/书签。建议用 list_pdf_info 了解页数后用 read_pdf_pages 逐段阅读，或用 search_pdf 按关键词定位。"
        lines = [f"📑 PDF 目录结构（共 {len(toc)} 条）："]
        for entry in toc:
            lines.append(f"  [{entry['page']}] {entry['title']}")
        return "\n".join(lines)
    except Exception as e:
        return f"提取 PDF 目录失败：{e}"


def read_pdf_pages(start_page: int = 1, end_page: int = 0, max_chars: int = 6000) -> str:
    """读取 PDF 指定页码范围。用于长文档分段阅读，避免一次加载全部内容。
    
    Args:
        start_page: 起始页码（从 1 开始），默认 1
        end_page: 结束页码（含），默认 0 表示只读 start_page 那一页
        max_chars: 最大返回字符数，默认 6000
    """
    try:
        from pypdf import PdfReader
    except ImportError:
        return "PDF 解析依赖 pypdf 未安装，请运行：pip install pypdf"

    try:
        _, data = _uploaded_file()
        reader = PdfReader(io.BytesIO(data))
        total = len(reader.pages)
        if start_page < 1:
            start_page = 1
        if end_page < start_page:
            end_page = start_page
        if start_page > total:
            return f"起始页 {start_page} 超出文档总页数 {total}。"

        end_page = min(end_page, total)
        pages = []
        for i in range(start_page - 1, end_page):
            text = reader.pages[i].extract_text() or "(该页无可提取文字)"
            pages.append(f"## 第 {i + 1} 页\n{text.strip()}")

        full = f"PDF 共 {total} 页，读取第 {start_page}-{end_page} 页：\n\n" + "\n\n".join(pages)
        if len(full) > max_chars:
            full = full[:max_chars] + f"\n\n…(已截断，原文 {len(full)} 字符。可缩小页面范围或用 search_pdf 定位关键内容)"
        return full
    except Exception as e:
        return f"读取 PDF 页面失败：{e}"


def search_pdf(query: str, context_chars: int = 500) -> str:
    """在 PDF 中搜索关键词，返回匹配页面及上下文。用于长文档快速定位关键信息。
    
    Args:
        query: 搜索关键词（支持简单文本匹配）
        context_chars: 每个匹配项返回的上下文字符数，默认 500
    """
    try:
        from pypdf import PdfReader
    except ImportError:
        return "PDF 解析依赖 pypdf 未安装，请运行：pip install pypdf"

    try:
        _, data = _uploaded_file()
        reader = PdfReader(io.BytesIO(data))
        total = len(reader.pages)
        results = []
        query_lower = query.lower()

        for i, page in enumerate(reader.pages):
            text = page.extract_text() or ""
            if not text:
                continue
            idx = text.lower().find(query_lower)
            if idx >= 0:
                start = max(0, idx - context_chars // 2)
                end = min(len(text), idx + len(query) + context_chars // 2)
                snippet = text[start:end]
                if start > 0:
                    snippet = "…" + snippet
                if end < len(text):
                    snippet = snippet + "…"
                results.append({
                    "page": i + 1,
                    "snippet": snippet,
                })

        if not results:
            return f"在 PDF 的 {total} 页中未找到关键词「{query}」。\n建议：①检查拼写；②尝试近义词；③用 list_pdf_info 了解文档结构再分段阅读。"

        lines = [f"🔍 搜索「{query}」——PDF 共 {total} 页，找到 {len(results)} 处匹配："]
        for r in results[:10]:
            lines.append(f"\n### 第 {r['page']} 页\n{r['snippet']}")
        if len(results) > 10:
            lines.append(f"\n…共 {len(results)} 处匹配，以上为前 10 处。可缩小搜索范围或指定页码范围进一步定位。")
        return "\n".join(lines)
    except Exception as e:
        return f"搜索 PDF 失败：{e}"
