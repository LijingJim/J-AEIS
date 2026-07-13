"""
=============================================================
  RAG 知识库 —— 文档向量化 + 语义检索
=============================================================
基于 ChromaDB + sentence-transformers（本地模型，零 API 成本）
- 上传文件自动切片、向量化、入库
- 对话中 LLM 可调用 search_knowledge_base 检索相关片段
=============================================================
"""
import hashlib
import os
import re
from pathlib import Path
from typing import List, Optional

# ChromaDB 持久化目录
CHROMA_DIR = Path(__file__).parent / "chroma_db"

# 文档集合名（按文件 hash 隔离，避免重复索引）
COLLECTION_PREFIX = "doc"

# 切片参数
CHUNK_SIZE = 600       # 每段字符数
CHUNK_OVERLAP = 100    # 段间重叠字符数


def _get_embedding_fn():
    """懒加载 ChromaDB 默认 embedding 函数（sentence-transformers all-MiniLM-L6-v2）。"""
    try:
        from chromadb.utils import embedding_functions
        return embedding_functions.DefaultEmbeddingFunction()
    except ImportError:
        raise ImportError(
            "请先安装 chromadb：pip install chromadb"
        )


def _file_hash(data: bytes) -> str:
    """计算文件内容 SHA256 前 12 位，用于去重。"""
    return hashlib.sha256(data).hexdigest()[:12]


def _chunk_text(text: str, chunk_size: int = CHUNK_SIZE, overlap: int = CHUNK_OVERLAP) -> List[str]:
    """将长文本按固定长度切片，保留重叠。"""
    if not text or not text.strip():
        return []
    text = text.strip()
    chunks = []
    start = 0
    while start < len(text):
        end = min(start + chunk_size, len(text))
        # 尽量在句号、换行处断开
        if end < len(text):
            for sep in ("\n\n", "\n", "。", "；", ". ", " "):
                pos = text.rfind(sep, start + chunk_size // 2, end)
                if pos > start:
                    end = pos + len(sep)
                    break
        chunk = text[start:end].strip()
        if chunk:
            chunks.append(chunk)
        start = end - overlap if end < len(text) else end
    return chunks


def _extract_text_from_bytes(data: bytes, filename: str) -> str:
    """从文件字节中提取纯文本。支持 txt/md/csv/json/pdf/docx/pptx/xlsx。"""
    ext = Path(filename).suffix.lower()

    # 纯文本
    if ext in {".txt", ".md", ".log", ".py", ".js", ".ts", ".sql", ".yaml", ".yml", ".ini", ".cfg"}:
        for enc in ("utf-8-sig", "utf-8", "gb18030", "gbk", "latin-1"):
            try:
                return data.decode(enc)
            except UnicodeDecodeError:
                continue
        return data.decode("utf-8", errors="replace")

    # CSV / JSON
    if ext == ".csv":
        try:
            text = data.decode("utf-8-sig")
            return text
        except Exception:
            pass

    if ext == ".json":
        try:
            import json
            return json.dumps(json.loads(data), ensure_ascii=False, indent=2)
        except Exception:
            return data.decode("utf-8", errors="replace")

    # PDF
    if ext == ".pdf":
        try:
            from pypdf import PdfReader
            reader = PdfReader(__import__("io").BytesIO(data))
            pages = []
            for page in reader.pages:
                t = page.extract_text()
                if t:
                    pages.append(t)
            return "\n\n".join(pages)
        except Exception:
            pass

    # DOCX
    if ext == ".docx":
        try:
            from docx import Document
            doc = Document(__import__("io").BytesIO(data))
            return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
        except Exception:
            pass

    # PPTX
    if ext == ".pptx":
        try:
            from pptx import Presentation
            prs = Presentation(__import__("io").BytesIO(data))
            texts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        texts.append(shape.text_frame.text)
            return "\n".join(texts)
        except Exception:
            pass

    # XLSX
    if ext in {".xlsx", ".xls"}:
        try:
            import pandas as pd
            df = pd.read_excel(__import__("io").BytesIO(data))
            return df.to_markdown(index=False) if hasattr(df, "to_markdown") else df.to_string(index=False)
        except Exception:
            pass

    # 兜底：尝试当文本读
    try:
        return data.decode("utf-8", errors="replace")
    except Exception:
        pass

    return ""


def index_file(data: bytes, filename: str, session_id: str = "") -> int:
    """
    将文件内容切片并存入向量库（按 session_id 隔离）。
    返回入库的切片数量。已索引过的文件（同 session_id + 同 hash）会跳过。
    """
    fhash = _file_hash(data)
    collection_name = f"{session_id}_{fhash}" if session_id else f"global_{fhash}"

    import chromadb
    client = chromadb.PersistentClient(path=str(CHROMA_DIR))

    # 检查是否已索引（count > 0 才视为有效索引）
    try:
        existing = client.get_collection(collection_name)
        if existing.count() > 0:
            return existing.count()
        client.delete_collection(collection_name)
    except Exception:
        pass

    text = _extract_text_from_bytes(data, filename)
    if not text or not text.strip():
        return 0

    chunks = _chunk_text(text)
    if not chunks:
        return 0

    embedding_fn = _get_embedding_fn()

    collection = client.create_collection(
        name=collection_name,
        embedding_function=embedding_fn,
        metadata={"hash": fhash, "filename": filename, "session_id": session_id},
    )

    ids = [f"{fhash}_{i}" for i in range(len(chunks))]
    metadatas = [{"filename": filename, "chunk_index": i, "hash": fhash, "session_id": session_id} for i in range(len(chunks))]

    BATCH = 100
    for i in range(0, len(chunks), BATCH):
        collection.add(
            ids=ids[i:i + BATCH],
            documents=chunks[i:i + BATCH],
            metadatas=metadatas[i:i + BATCH],
        )

    return len(chunks)


def search_knowledge_base(query: str, session_id: str = "", top_k: int = 5) -> List[dict]:
    """
    在当前会话的已索引文件中语义检索相关片段。
    session_id 为空时搜索全局库。
    """
    embedding_fn = _get_embedding_fn()
    import chromadb
    client = chromadb.PersistentClient(path=str(CHROMA_DIR))

    collections = client.list_collections()
    if not collections:
        return []

    # 只搜索当前 session 的 collection
    prefix = f"{session_id}_" if session_id else "global_"
    my_collections = [c for c in collections if c.name.startswith(prefix)]

    all_results = []
    for col in my_collections:
        try:
            results = col.query(query_texts=[query], n_results=top_k)
            docs = results.get("documents", [[]])[0]
            metas = results.get("metadatas", [[]])[0]
            distances = results.get("distances", [[]])[0]
            for j, doc in enumerate(docs):
                meta = metas[j] if j < len(metas) else {}
                all_results.append({
                    "filename": meta.get("filename", "?"),
                    "chunk_index": meta.get("chunk_index", -1),
                    "content": doc,
                    "hash": meta.get("hash", "?"),
                    "score": round(1 - distances[j], 4) if j < len(distances) else 0,
                })
        except Exception:
            continue

    seen = set()
    unique = []
    for r in sorted(all_results, key=lambda x: x["score"], reverse=True):
        key = r["content"][:80]
        if key not in seen:
            seen.add(key)
            unique.append(r)
        if len(unique) >= top_k:
            break

    return unique


def clear_knowledge_base(session_id: str = ""):
    """清空指定会话的向量索引。session_id 为空时清空全部。"""
    import chromadb
    client = chromadb.PersistentClient(path=str(CHROMA_DIR))
    prefix = f"{session_id}_" if session_id else ""
    for col in client.list_collections():
        if not prefix or col.name.startswith(prefix):
            client.delete_collection(col.name)


def get_indexed_files(session_id: str = "") -> List[dict]:
    """返回当前会话已索引的文件列表。session_id 为空时返回全部。"""
    import chromadb
    client = chromadb.PersistentClient(path=str(CHROMA_DIR))
    prefix = f"{session_id}_" if session_id else ""
    files = []
    for col in client.list_collections():
        if prefix and not col.name.startswith(prefix):
            continue
        meta = col.metadata or {}
        files.append({
            "hash": meta.get("hash", "?"),
            "filename": meta.get("filename", "?"),
            "chunks": col.count(),
        })
    return files
